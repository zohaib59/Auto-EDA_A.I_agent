import streamlit as st
import pandas as pd
import plotly.express as px
import seaborn as sns
import matplotlib.pyplot as plt
import io
from fpdf import FPDF
from pptx import Presentation
from pptx.util import Inches
from dotenv import load_dotenv
from langchain_openai import ChatOpenAI
from langchain.chains import RetrievalQA
from langchain.document_loaders import CSVLoader
from langchain.vectorstores import FAISS
from langchain.embeddings import OpenAIEmbeddings

load_dotenv()
st.set_page_config(layout="wide")
st.title("⚡ Smart AI-Powered EDA App")
st.markdown("Upload a CSV file to perform automated EDA, ask questions, and export to PDF/PPT.")

# Session Control Buttons
col_kill, col_clear = st.columns(2)
if col_kill.button("🔚 End Session"):
    st.session_state.clear()
    st.success("Session Ended")
if col_clear.button("🗑️ Clear Chat History"):
    st.session_state.pop("chat_history", None)
    st.success("Chat History Cleared")

# Dark Mode
if st.toggle("🌙 Toggle Dark Mode"):
    st.markdown("""<style>body, .stApp {background-color: #1E1E1E; color: white;}</style>""", unsafe_allow_html=True)

# File Upload
uploaded_file = st.file_uploader("📁 Upload CSV", type=["csv"])

@st.cache_data
def load_data(file):
    df = pd.read_csv(file)
    df.columns = df.columns.str.strip()
    df = df.apply(lambda x: x.str.strip() if x.dtype == "object" else x)
    df = df.drop_duplicates()
    df = df.fillna(df.mean(numeric_only=True))
    return df

if uploaded_file:
    df = load_data(uploaded_file)

    st.subheader("📊 Dataset Preview")
    st.dataframe(df.head())

    st.subheader("📌 Dataset Summary")
    st.write("**Shape:**", df.shape)
    st.write("**Missing Values:**", df.isnull().sum().sum())
    st.write("**Duplicates Removed:**", df.duplicated().sum())
    st.dataframe(df.describe())

    numeric_cols = df.select_dtypes(include='number').columns.tolist()
    categorical_cols = df.select_dtypes(include='object').columns.tolist()
    target_col = df.columns[-1] if df.columns[-1] in numeric_cols + categorical_cols else None
    feature_cols = [col for col in df.columns if col != target_col]

    st.subheader("🎯 KPI Metrics")
    for col in numeric_cols:
        st.metric(label=f"{col} Sum", value=f"{df[col].sum():,.2f}")
        st.metric(label=f"{col} Avg", value=f"{df[col].mean():,.2f}")
        st.metric(label=f"{col} Max", value=f"{df[col].max():,.2f}")
        st.metric(label=f"{col} Min", value=f"{df[col].min():,.2f}")

    st.subheader("📊 Top 5 Values (All Columns)")
    for col in numeric_cols:
        top5_sum = df.groupby(col).size().reset_index(name='Count').nlargest(5, 'Count')
        st.write(f"Top 5 by count - {col}")
        st.dataframe(top5_sum)
    for col in categorical_cols:
        top5_cat = df[col].value_counts().nlargest(5).reset_index()
        top5_cat.columns = [col, 'Count']
        st.write(f"Top 5 categories - {col}")
        st.dataframe(top5_cat)

    st.subheader("📊 Top 5 / Bottom 5 as Bar Charts")
    for col in numeric_cols[:3]:
        top5 = df[[col]].nlargest(5, col)
        bottom5 = df[[col]].nsmallest(5, col)
        st.plotly_chart(px.bar(top5, y=col, title=f"Top 5 {col}"), use_container_width=True)
        st.plotly_chart(px.bar(bottom5, y=col, title=f"Bottom 5 {col}"), use_container_width=True)

    st.subheader("🥧 Enhanced Pie Charts")
    for col in categorical_cols[:3]:
        vc = df[col].value_counts().nlargest(5)
        fig = px.pie(values=vc.values, names=vc.index, title=f"Top 5 {col}")
        fig.update_traces(textinfo='percent+value')
        st.plotly_chart(fig, use_container_width=True)

    st.subheader("📈 Feature vs Target Scatter Plots")
    if target_col and target_col in numeric_cols:
        for col in feature_cols[:3]:
            if col in numeric_cols:
                fig = px.scatter(df, x=col, y=target_col, trendline="ols", title=f"{col} vs {target_col}")
                st.plotly_chart(fig, use_container_width=True)

    st.subheader("📊 Summary Stats Table")
    stats = df.describe().T[['mean', 'min', 'max']]
    st.dataframe(stats)

    st.subheader("📉 Historical Trends")
    date_cols = df.select_dtypes(include='datetime64').columns.tolist()
    if not date_cols:
        try:
            df['parsed_date'] = pd.to_datetime(df.iloc[:, 0], errors='coerce')
            if df['parsed_date'].notna().sum() > 0:
                date_cols.append('parsed_date')
        except:
            pass
    if date_cols:
        df = df.dropna(subset=date_cols)
        df = df.sort_values(by=date_cols[0])
        for col in numeric_cols[:2]:
            fig = px.line(df, x=date_cols[0], y=col, title=f"Trend of {col} Over Time")
            st.plotly_chart(fig, use_container_width=True)

    st.subheader("📦 Outlier Detection")
    for col in numeric_cols[:3]:
        st.plotly_chart(px.box(df, y=col, title=f"Box Plot: {col}"), use_container_width=True)

    st.subheader("📊 Categorical vs Numerical Comparison")
    for cat in categorical_cols[:2]:
        for num in numeric_cols[:2]:
            st.plotly_chart(px.box(df, x=cat, y=num, title=f"{num} by {cat}"), use_container_width=True)

    st.subheader("💬 Ask Your Data Anything")
    if "qa_chain" not in st.session_state:
        with open("temp.csv", "wb") as f:
            f.write(uploaded_file.getbuffer())
        loader = CSVLoader(file_path="temp.csv")
        docs = loader.load()
        vectorstore = FAISS.from_documents(docs, OpenAIEmbeddings())
        qa_chain = RetrievalQA.from_chain_type(llm=ChatOpenAI(model="gpt-3.5-turbo"), retriever=vectorstore.as_retriever())
        st.session_state.qa_chain = qa_chain

    query = st.text_input("Ask a question about the dataset:")
    if query:
        st.write("🤖", st.session_state.qa_chain.run(query))

    st.subheader("📤 Export Summary")
    def generate_pdf():
        pdf = FPDF()
        pdf.add_page()
        pdf.set_font("Arial", size=12)
        pdf.cell(200, 10, txt="EDA Summary Report", ln=True, align='C')
        for col in df.columns:
            nulls = df[col].isnull().sum()
            dtype = df[col].dtype
            pdf.cell(200, 10, txt=f"{col} - {dtype}, Nulls: {nulls}", ln=True)
        return pdf.output(dest="S").encode("latin1")

    def generate_ppt():
        prs = Presentation()
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = "EDA Highlights"
        for idx, col in enumerate(df.columns[:5]):
            txBox = slide.shapes.add_textbox(Inches(1), Inches(1 + idx * 0.6), Inches(8), Inches(0.5))
            tf = txBox.text_frame
            tf.text = f"{col}: {df[col].dtype}, Nulls: {df[col].isnull().sum()}"
        return prs

    col1, col2 = st.columns(2)
    with col1:
        if st.download_button("📄 Download PDF", generate_pdf(), "eda_report.pdf", mime="application/pdf"):
            st.success("PDF Downloaded")

    with col2:
        ppt_file = io.BytesIO()
        generate_ppt().save(ppt_file)
        ppt_file.seek(0)
        if st.download_button("📊 Download PPT", ppt_file, "eda_slides.pptx", mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"):
            st.success("PPT Downloaded")
